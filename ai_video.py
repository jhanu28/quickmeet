#!/usr/bin/env python
# ai_video.py

import os
import sys
import types
import concurrent.futures
from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Inches
from pdf2image import convert_from_path
from PIL import Image

# ========== 🐵 MONKEY PATCH MOVIEPY ==========
def dummy_concatenate_videoclips(clips, method="compose"):
    """A dummy function to replace the missing concatenate_videoclips function."""
    from moviepy.editor import VideoClip
    cum_durations = [0]
    for clip in clips:
        cum_durations.append(cum_durations[-1] + clip.duration)
    total_duration = cum_durations[-1]

    def make_frame(t):
        for i in range(len(clips)):
            if cum_durations[i] <= t < cum_durations[i+1]:
                return clips[i].get_frame(t - cum_durations[i])
        return clips[-1].get_frame(clips[-1].duration)

    return VideoClip(make_frame, duration=total_duration)

# Manually create a fake `moviepy.video.compositing.concatenate` module
fake_module = types.ModuleType("moviepy.video.compositing.concatenate")
fake_module.concatenate_videoclips = dummy_concatenate_videoclips
sys.modules["moviepy.video.compositing.concatenate"] = fake_module

# Now safely import moviepy
from moviepy.editor import VideoClip, ImageClip, CompositeVideoClip, VideoFileClip

print("✅ MoviePy imported successfully without errors.")

# ========== 🛠 Load AWS Credentials ==========
load_dotenv()
AWS_ACCESS_KEY = os.getenv("AWS_ACCESS_KEY_ID")
AWS_SECRET_KEY = os.getenv("AWS_SECRET_ACCESS_KEY")
S3_BUCKET = "quickmeet-files"

# ========== 📁 Ensure Required Directories Exist ==========
os.makedirs("slides", exist_ok=True)

# ========== 🔽 AWS S3 FUNCTIONS ==========
def download_from_s3(s3_key, local_filename):
    """Download a file from S3 with error handling."""
    import boto3  # Import here to avoid issues if not installed globally
    s3 = boto3.client(
        's3',
        region_name='us-east-1',
        aws_access_key_id=AWS_ACCESS_KEY,
        aws_secret_access_key=AWS_SECRET_KEY
    )
    try:
        s3.download_file(S3_BUCKET, s3_key, local_filename)
        print(f"✅ Downloaded {s3_key} from S3.")
        return True
    except Exception as e:
        print(f"❌ Error downloading {s3_key}: {e}")
        return False

# ========== 📝 GENERATE PPTX FROM TEXT ==========
def create_ppt_from_text(summary_file, action_items_file, ppt_filename):
    """Generate a PowerPoint presentation from text files."""
    prs = Presentation()

    def add_slides(title, content):
        max_chars = 300
        words = content.split()
        slides_text = []
        current_slide = ""
        for word in words:
            if len(current_slide) + len(word) + 1 <= max_chars:
                current_slide += (" " + word) if current_slide else word
            else:
                slides_text.append(current_slide)
                current_slide = word
        if current_slide:
            slides_text.append(current_slide)
        for chunk in slides_text:
            slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(slide_layout)
            title_shape = slide.shapes.title
            if title_shape:
                title_shape.text = title
            text_box = slide.shapes.add_textbox(0, Inches(1.5), prs.slide_width, prs.slide_height - Inches(1.5))
            text_box.text_frame.text = chunk

    with open(summary_file, "r") as file:
        summary = file.read().strip()
    with open(action_items_file, "r") as file:
        action_items = file.read().strip()

    add_slides("Meeting Summary", summary)
    add_slides("Action Items", action_items)

    prs.save(ppt_filename)
    print(f"✅ Created PowerPoint: {ppt_filename}")

# ========== 🖼 CONVERT PPTX TO IMAGES ==========
def convert_ppt_to_images_parallel(ppt_filename, output_folder):
    """Convert PowerPoint slides to images using multiprocessing."""
    pdf_filename = ppt_filename.replace(".pptx", ".pdf")
    os.system(f'soffice --headless --convert-to pdf --outdir . "{ppt_filename}"')
    if not os.path.exists(pdf_filename):
        print("❌ Error: PDF conversion failed.")
        return []
    images = convert_from_path(pdf_filename)
    os.makedirs(output_folder, exist_ok=True)
    
    def save_image(index, image):
        image_path = os.path.join(output_folder, f"slide_{index+1}.png")
        image.save(image_path, "PNG")
        print(f"✅ Processed {image_path}")
    
    with concurrent.futures.ThreadPoolExecutor() as executor:
        executor.map(save_image, range(len(images)), images)
    
    return [os.path.join(output_folder, f"slide_{i+1}.png") for i in range(len(images))]

# ========== 🎥 MERGE VIDEO & SLIDES ==========
def merge_video_with_slides(avatar_video, slides_folder, output_video):
    """Merge the AI avatar video with PPT slides as a background."""
    slide_images = sorted(
        [os.path.join(slides_folder, img) for img in os.listdir(slides_folder) if img.endswith(".png")]
    )
    if not slide_images:
        print("❌ No slides found for merging.")
        return

    avatar_clip = VideoFileClip(avatar_video)
    total_duration = avatar_clip.duration
    num_slides = len(slide_images)
    slide_duration = total_duration / num_slides

    slide_clips = [ImageClip(img).set_duration(slide_duration) for img in slide_images]

    # Import patched concatenate_videoclips
    from moviepy.video.compositing.concatenate import concatenate_videoclips

    background = concatenate_videoclips(slide_clips, method="compose")

    final_video = CompositeVideoClip([background, avatar_clip.set_position(("center", "bottom"))], size=(1280, 720))
    final_video.write_videofile(output_video, codec="libx264", fps=24)
    print(f"✅ Final video saved as {output_video}")

# ========== 🚀 MAIN EXECUTION ==========
if __name__ == "__main__":
    summary_file = "summary.txt"
    action_items_file = "action_items.txt"
    ppt_filename = "meeting_summary.pptx"
    slides_folder = "slides"
    avatar_video = "avatar_video.mp4"
    final_video = "final_meeting_video.mp4"

    # Download required files from S3
    download_from_s3("summaries/summary.txt", summary_file)
    download_from_s3("summaries/action_items.txt", action_items_file)
    download_from_s3("videos/avatar_video.mp4", avatar_video)

    # Step 1: Create PowerPoint from text files
    create_ppt_from_text(summary_file, action_items_file, ppt_filename)

    # Step 2: Convert PPTX to images
    convert_ppt_to_images_parallel(ppt_filename, slides_folder)

    # Step 3: Merge video with slides
    merge_video_with_slides(avatar_video, slides_folder, final_video)
