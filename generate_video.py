#!/usr/bin/env python
# generate_video.py

# ===== Monkey Patch Section =====
import sys
import types

# --- Patch for missing moviepy.video.compositing.concatenate ---
def dummy_concatenate_videoclips(clips, method="compose"):
    """
    A basic concatenation function that stitches video clips sequentially.
    Mimics MoviePy's concatenate_videoclips.
    """
    from moviepy.video.VideoClip import VideoClip
    cum_durations = [0]
    for clip in clips:
        cum_durations.append(cum_durations[-1] + clip.duration)
    total_duration = cum_durations[-1]

    def make_frame(t):
        for i in range(len(clips)):
            if cum_durations[i] <= t < cum_durations[i+1]:
                return clips[i].get_frame(t - cum_durations[i])
        return clips[-1].get_frame(clips[-1].duration)

    return VideoClip(make_frame, duration=total_duration)

dummy_concat_module = types.ModuleType("moviepy.video.compositing.concatenate")
dummy_concat_module.concatenate_videoclips = dummy_concatenate_videoclips
sys.modules["moviepy.video.compositing.concatenate"] = dummy_concat_module

# --- Patch for missing "all" attribute in moviepy.audio.fx ---
if "moviepy.audio.fx" not in sys.modules:
    dummy_audio_fx = types.ModuleType("moviepy.audio.fx")
    sys.modules["moviepy.audio.fx"] = dummy_audio_fx
else:
    dummy_audio_fx = sys.modules["moviepy.audio.fx"]
dummy_audio_fx.all = []

# --- Patch for missing moviepy.video.fx.fadein ---
dummy_fadein_module = types.ModuleType("moviepy.video.fx.fadein")
def dummy_fadein(clip, duration, initial_color=None):
    return clip
dummy_fadein_module.fadein = dummy_fadein
sys.modules["moviepy.video.fx.fadein"] = dummy_fadein_module

# --- Patch for missing moviepy.video.fx.fadeout ---
dummy_fadeout_module = types.ModuleType("moviepy.video.fx.fadeout")
def dummy_fadeout(clip, duration, final_color=None):
    return clip
dummy_fadeout_module.fadeout = dummy_fadeout
sys.modules["moviepy.video.fx.fadeout"] = dummy_fadeout_module

def resize_clip(clip, width=None, height=None):
    if width and height:
        return clip.resized((width, height))
    elif width:
        return clip.resized(width=width)
    elif height:
        return clip.resized(height=height)
    return clip 

# ===== End Monkey Patch Section =====

import os
import boto3
import requests
import json
import concurrent.futures
from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Inches
from pdf2image import convert_from_path
from PIL import Image
from moviepy.video.io.VideoFileClip import VideoFileClip
from moviepy.audio.io.AudioFileClip import AudioFileClip
# Import ImageClip and CompositeVideoClip directly from their modules
from moviepy.video.VideoClip import ImageClip
from moviepy.video.compositing.CompositeVideoClip import CompositeVideoClip
# Import our dummy-patched concatenate function
from moviepy.video.compositing.concatenate import concatenate_videoclips
import time
from moviepy.editor import VideoFileClip, vfx

# Load AWS credentials from .env
load_dotenv()
AWS_ACCESS_KEY = os.getenv("AWS_ACCESS_KEY_ID")
AWS_SECRET_KEY = os.getenv("AWS_SECRET_ACCESS_KEY")
S3_BUCKET = "quickmeet-files"
HEYGEN_API_KEY = "MjNhYTNhYzIzZjU0NDExYTk2NjU0NzhjOTM1YjFlYzctMTczOTg3MTg5MQ=="

# Ensure required directories exist
os.makedirs("slides", exist_ok=True)

# -------------------- AWS S3 FUNCTIONS --------------------
def upload_to_s3(local_filename, s3_key):
    """Upload a file to S3 with retries."""
    s3 = boto3.client(
        's3',
        region_name='us-east-1',
        aws_access_key_id=AWS_ACCESS_KEY,
        aws_secret_access_key=AWS_SECRET_KEY
    )
    for attempt in range(3):
        try:
            s3.upload_file(local_filename, S3_BUCKET, s3_key)
            print(f"✅ Uploaded {local_filename} to S3 as {s3_key}")
            return True
        except Exception as e:
            print(f"⚠️ Upload failed (Attempt {attempt + 1}): {e}")
            time.sleep(2)
    print(f"❌ Upload failed after multiple attempts: {local_filename}")
    return False

def download_from_s3(s3_key, local_filename):
    """Download a file from S3 with error handling."""
    s3 = boto3.client(
        's3',
        region_name='us-east-1',
        aws_access_key_id=AWS_ACCESS_KEY,
        aws_secret_access_key=AWS_SECRET_KEY
    )
    try:
        s3.download_file(S3_BUCKET, s3_key, local_filename)
        print(f"✅ Downloaded {s3_key} from S3.")
        return True
    except Exception as e:
        print(f"❌ Error downloading {s3_key}: {e}")
        return False

# -------------------- GENERATE PPTX --------------------
def create_ppt_from_text(summary_file, action_items_file, ppt_filename):
    """
    Generate a PowerPoint presentation from text files.
    This function splits text into multiple slides if necessary.
    For simplicity, it creates one slide per file unless text exceeds 300 characters.
    """
    prs = Presentation()
    def add_slides(title, content):
        max_chars = 300  # Adjust as needed
        words = content.split()
        slides_text = []
        current_slide = ""
        for word in words:
            if len(current_slide) + len(word) + 1 <= max_chars:
                current_slide += (" " + word) if current_slide else word
            else:
                slides_text.append(current_slide)
                current_slide = word
        if current_slide:
            slides_text.append(current_slide)
        for chunk in slides_text:
            slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(slide_layout)
            title_shape = slide.shapes.title
            if title_shape:
                title_shape.text = title
            else:
                title_shape = slide.shapes.add_textbox(0, 0, prs.slide_width, Inches(1))
                title_shape.text = title
            text_box = slide.shapes.add_textbox(0, Inches(1.5), prs.slide_width, prs.slide_height - Inches(1.5))
            text_box.text_frame.text = chunk

    with open(summary_file, "r") as file:
        summary = file.read().strip()
    with open(action_items_file, "r") as file:
        action_items = file.read().strip()

    add_slides("Meeting Summary", summary)
    add_slides("Action Items", action_items)

    prs.save(ppt_filename)
    print(f"✅ Created PowerPoint: {ppt_filename}")
    upload_to_s3(ppt_filename, ppt_filename)

# -------------------- CONVERT PPTX TO IMAGES --------------------
def convert_ppt_to_images_parallel(ppt_filename, output_folder):
    """Convert PowerPoint slides to images using multiprocessing."""
    pdf_filename = ppt_filename.replace(".pptx", ".pdf")
    os.system(f'soffice --headless --convert-to pdf --outdir . "{ppt_filename}"')
    if not os.path.exists(pdf_filename):
        print("❌ Error: PDF conversion failed.")
        return []
    images = convert_from_path(pdf_filename)
    os.makedirs(output_folder, exist_ok=True)
    def save_image(index, image):
        image_path = os.path.join(output_folder, f"slide_{index+1}.png")
        image.save(image_path, "PNG")
        upload_to_s3(image_path, image_path)
        print(f"✅ Processed {image_path}")
    with concurrent.futures.ThreadPoolExecutor() as executor:
        executor.map(save_image, range(len(images)), images)
    return [os.path.join(output_folder, f"slide_{i+1}.png") for i in range(len(images))]

# -------------------- GENERATE VIDEO FROM HEYGEN --------------------
def generate_ai_video(text_script, output_video):
    """
    Generate an AI video using a HeyGen avatar with your custom configuration.
    The avatar reads the full text (summary & action items) and the function waits
    until the video is fully processed before returning.
    """
    url = "https://api.heygen.com/v2/video/generate"
    headers = {
        "X-Api-Key": HEYGEN_API_KEY,
        "Content-Type": "application/json"
    }
    payload = {
        "video_inputs": [
            {
                "character": {
                    "type": "avatar",
                    "avatar_id": "Miyu_sitting_sofabusiness_front",
                    "avatar_style": "normal",
                    "offset": {"x": 0, "y": 0.28},
                    "scale": 1.33
                },
                "voice": {
                    "type": "text",
                    "input_text": text_script,
                    "voice_id": "2d5b0e6cf36f460aa7fc47e3eee4ba54"
                },
                "background": {
                    "type": "color",
                    "value": "#ffffff"
                }
            }
        ],
        "dimension": {
            "width": 1280,
            "height": 720
        }
    }
    
    import time, json, requests
    for attempt in range(3):
        response = requests.post(url, headers=headers, json=payload)
        try:
            response.raise_for_status()
            data = response.json()
        except Exception as e:
            print("❌ Failed to generate video. Error:", e)
            time.sleep(2)
            continue

        print("DEBUG: API response (attempt {}):".format(attempt + 1))
        print(data)
        
        video_id = data.get("data", {}).get("video_id")
        if not video_id:
            print("❌ Video ID not found in response.")
            time.sleep(2)
            continue
        print(f"✅ Video created with video_id: {video_id}")
        
        # Poll for video status until it's completed.
        status_url = f"https://api.heygen.com/v1/video_status.get?video_id={video_id}"
        headers_status = {
            "X-Api-Key": HEYGEN_API_KEY,
            "Accept": "application/json"
        }
        timeout = 600  # Extended timeout: wait up to 10 minutes
        interval = 5
        start_time = time.time()
        video_url = None
        while time.time() - start_time < timeout:
            status_response = requests.get(status_url, headers=headers_status)
            try:
                status_data = status_response.json()
            except json.JSONDecodeError:
                print("❌ Failed to parse JSON from status response.")
                time.sleep(interval)
                continue
            
            # Extract status from nested "data" field
            status = status_data.get("data", {}).get("status")
            print(f"Video status: {status}")
            if status == "completed":
                video_url = status_data.get("data", {}).get("video_url")
                break
            elif status == "failed":
                print("❌ Video generation failed.")
                return False
            time.sleep(interval)
        
        if not video_url:
            print("❌ Video URL not found after polling.")
            time.sleep(2)
            continue
        
        video_response = requests.get(video_url)
        if video_response.status_code == 200:
            with open(output_video, "wb") as f:
                f.write(video_response.content)
            print(f"✅ AI Video saved to {output_video}")
            return True
        else:
            print(f"❌ Error downloading video from URL: {video_url}")
        time.sleep(2)
    
    print("❌ Failed to generate AI video after multiple attempts.")
    return False

# -------------------- MERGE VIDEO & SLIDES WITH SYNCED SLIDES --------------------
def merge_video_with_slides(avatar_video, slides_folder, output_video):
    """
    Merge the AI avatar video with PPT slides as a background.
    The PPT slides are displayed in sequence so that each slide covers
    a segment of the avatar video. The duration for each slide is determined
    by evenly dividing the avatar video's total duration.
    """
    slide_images = sorted(
        [os.path.join(slides_folder, img) for img in os.listdir(slides_folder) if img.endswith(".png")]
    )
    if not slide_images:
        print("❌ No slides found for merging.")
        return

    # Load the full avatar video to get its duration.
    avatar_clip = VideoFileClip(avatar_video)
    total_duration = avatar_clip.duration
    num_slides = len(slide_images)
    duration_per_slide = total_duration / num_slides

    # Create an ImageClip for each slide with the calculated duration.
    clips = [ImageClip(img).with_duration(duration_per_slide) for img in slide_images]
    background_clip = concatenate_videoclips(clips, method="compose")
    
    # Resize avatar clip to a smaller width and position it at bottom-right.
    avatar_overlay = resize_clip(avatar_clip, width=300).set_position(("right", "bottom"))

    # Composite the background (PPT slides) and the avatar overlay.
    final_clip = CompositeVideoClip([background_clip, avatar_overlay])
    
    final_clip.write_videofile(output_video, codec="libx264", fps=24)
    print(f"✅ Final Video Created: {output_video}")
    upload_to_s3(output_video, output_video)

# -------------------- MAIN FUNCTION --------------------
def main():
    summary_txt = "summary.txt"
    action_items_txt = "action_items.txt"
    ppt_filename = "meeting_summary.pptx"
    slides_folder = "slides"
    avatar_video = "avatar_video.mp4"
    final_video = "final_meeting_video.mp4"

    if not download_from_s3("summary.txt", summary_txt) or not download_from_s3("action_items.txt", action_items_txt):
        print("❌ Failed to download required files. Exiting.")
        return

    create_ppt_from_text(summary_txt, action_items_txt, ppt_filename)
    slide_images = convert_ppt_to_images_parallel(ppt_filename, slides_folder)
    if not slide_images:
        return

    with open(summary_txt, "r") as file:
        text_script = file.read()

    if not generate_ai_video(text_script, avatar_video):
        return

    merge_video_with_slides(avatar_video, slides_folder, final_video)
    print("🎉 Process Completed! Final video is in S3.")

if __name__ == "__main__":
    main()
